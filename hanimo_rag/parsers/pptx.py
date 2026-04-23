"""PPTX parser using python-pptx."""
from __future__ import annotations

from hanimo_rag.parsers.base import ParsedDocument, ParserBase


def _extract_slide_text(slide) -> str:
    """Extract all text from a slide's shapes."""
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                parts.append(line)
    return "\n".join(parts)


def _extract_slide_notes(slide) -> str:
    """Extract speaker notes from a slide."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return tf.text.strip() if tf else ""
    except Exception:
        return ""


class PptxParser(ParserBase):
    """Parse PPTX files extracting text and notes from each slide."""

    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]

    def parse(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError as e:
            return ParsedDocument(
                text="",
                metadata={"error": f"Missing dependency: {e}"},
            )

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            return ParsedDocument(text="", metadata={"error": str(exc)})

        slides = prs.slides
        metadata: dict = {"slide_count": len(slides)}

        # Try to get presentation title from first slide
        try:
            first_slide = slides[0]
            for shape in first_slide.shapes:
                if shape.has_text_frame and shape.shape_type == 13:  # MSO_SHAPE_TYPE.TITLE
                    metadata["title"] = shape.text_frame.text.strip()
                    break
            if "title" not in metadata:
                # Fallback: first non-empty text shape
                for shape in first_slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            metadata["title"] = t
                            break
        except (IndexError, Exception):
            pass

        pages: list[str] = []
        for i, slide in enumerate(slides, start=1):
            slide_text = _extract_slide_text(slide)
            notes = _extract_slide_notes(slide)

            parts = [f"## Slide {i}"]
            if slide_text:
                parts.append(slide_text)
            if notes:
                parts.append(f"**Notes:** {notes}")

            pages.append("\n\n".join(parts))

        full_text = "\n\n".join(pages)
        return ParsedDocument(text=full_text, metadata=metadata, pages=pages)
